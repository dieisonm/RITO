#!/usr/bin/env python3
from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable, List, Sequence, Tuple

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch, mm
from reportlab.pdfgen import canvas
from reportlab.platypus import Image, Paragraph, Preformatted, SimpleDocTemplate, Spacer


ROOT = Path(__file__).resolve().parent.parent
LOGO = ROOT / "assets" / "brand" / "logos" / "site" / "rito_sistemas_wordmark_01.png"
MONOGRAM = ROOT / "assets" / "brand" / "logos" / "site" / "rito_monogram_r_01.png"
OUT_ROOT = ROOT / "assets" / "business-kit"
DOCX_DIR = OUT_ROOT / "editable-docx"
PDF_DIR = OUT_ROOT / "pdf"
PPT_DIR = OUT_ROOT / "presentation"

BG_HEX = "#f4f0e9"
BG_LIGHT_HEX = "#fbf7f1"
SURFACE_HEX = "#fffdf9"
INK_HEX = "#152733"
INK_SOFT_HEX = "#5d6b72"
BRAND_HEX = "#173847"
BRAND_STRONG_HEX = "#0d2430"
BRAND_MID_HEX = "#315161"
ACCENT_HEX = "#b89163"
LINE_HEX = "#d7d1ca"


DOCX_SOURCES: List[Tuple[str, Path]] = [
    ("rito-first-contact-playbook.docx", ROOT / "docs/sales/client-ready/first-contact-playbook.md"),
    ("rito-client-briefing-form.docx", ROOT / "docs/sales/client-ready/client-briefing-form-final.md"),
    ("rito-commercial-proposal-template.docx", ROOT / "docs/sales/client-ready/commercial-proposal-final-template.md"),
    ("rito-estimate-template.docx", ROOT / "docs/sales/client-ready/estimate-final-template.md"),
    ("rito-contract-template.docx", ROOT / "docs/legal/contract-template.md"),
    ("rito-follow-up-sequence.docx", ROOT / "docs/sales/client-ready/follow-up-sequence-final.md"),
    ("rito-proposal-send-checklist.docx", ROOT / "docs/sales/client-ready/proposal-send-checklist.md"),
    ("rito-pricing-guide.docx", ROOT / "docs/sales/commercial/pricing-guide.md"),
]

PDF_SOURCES: List[Tuple[str, Path]] = [
    ("rito-branding-guide.pdf", ROOT / "docs/brand/institutional/mini-brand-guide.md"),
    ("rito-software-brand-system.pdf", ROOT / "docs/brand/institutional/software-ui-brand-system.md"),
    ("rito-company-onepager.pdf", ROOT / "docs/sales/materials/rito-company-onepager-v1.md"),
]

PRESENTATION_MD = ROOT / "docs/sales/materials/rito-company-presentation-v1.md"
PRESENTATION_PPTX = PPT_DIR / "rito-company-presentation.pptx"
PRESENTATION_PDF = PDF_DIR / "rito-company-presentation.pdf"

LABEL_LINES = {
    "texto de apoio:",
    "pontos de apoio:",
    "itens de exemplo:",
    "fluxo sugerido:",
    "mensagem de apoio:",
    "exemplos:",
    "benefícios principais:",
    "beneficios principais:",
    "perfil ideal:",
    "fluxo comercial:",
    "chamada final:",
}


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def docx_rgb(hex_value: str) -> DocxRGBColor:
    value = hex_value.lstrip("#")
    return DocxRGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def ensure_dirs() -> None:
    DOCX_DIR.mkdir(parents=True, exist_ok=True)
    PDF_DIR.mkdir(parents=True, exist_ok=True)
    PPT_DIR.mkdir(parents=True, exist_ok=True)


def existing_sources(sources: Sequence[Tuple[str, Path]]) -> Iterable[Tuple[str, Path]]:
    for filename, source in sources:
        if source.exists():
            yield filename, source


def clean_inline(text: str) -> str:
    text = re.sub(r"\[(.*?)\]\([^)]+\)", r"\1", text)
    text = text.replace("**", "").replace("__", "").replace("`", "")
    return text.strip()


def parse_markdown(text: str):
    blocks = []
    paragraph_lines: List[str] = []
    code_lines: List[str] = []
    in_code = False

    def flush_paragraph() -> None:
        if paragraph_lines:
            paragraph = " ".join(clean_inline(line.strip()) for line in paragraph_lines if line.strip())
            if paragraph:
                blocks.append(("paragraph", paragraph))
            paragraph_lines.clear()

    def flush_code() -> None:
        if code_lines:
            blocks.append(("code", "\n".join(code_lines)))
            code_lines.clear()

    for raw_line in text.splitlines():
        line = raw_line.rstrip("\n")
        stripped = line.strip()

        if stripped.startswith("```"):
            flush_paragraph()
            if in_code:
                flush_code()
                in_code = False
            else:
                in_code = True
            continue

        if in_code:
            code_lines.append(line)
            continue

        if not stripped:
            flush_paragraph()
            continue

        heading_match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        bullet_match = re.match(r"^[-*]\s+(.*)$", stripped)
        number_match = re.match(r"^(\d+)\.\s+(.*)$", stripped)
        quote_match = re.match(r"^>\s+(.*)$", stripped)

        if heading_match:
            flush_paragraph()
            blocks.append(("heading", len(heading_match.group(1)), clean_inline(heading_match.group(2))))
        elif bullet_match:
            flush_paragraph()
            blocks.append(("bullet", clean_inline(bullet_match.group(1))))
        elif number_match:
            flush_paragraph()
            blocks.append(("number", clean_inline(number_match.group(2))))
        elif quote_match:
            flush_paragraph()
            blocks.append(("quote", clean_inline(quote_match.group(1))))
        elif stripped.startswith("|"):
            flush_paragraph()
            blocks.append(("paragraph", clean_inline(stripped)))
        else:
            paragraph_lines.append(line)

    flush_paragraph()
    if in_code:
        flush_code()

    return blocks


def document_title(blocks, fallback: str) -> str:
    for block in blocks:
        if block[0] == "heading":
            return block[2]
    return fallback


def style_docx_run(run, *, font: str, size: float, bold: bool = False, italic: bool = False, color: str = INK_HEX) -> None:
    run.font.name = font
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = docx_rgb(color)


def setup_docx_document(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(0.6)
    section.bottom_margin = Inches(0.6)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    normal = doc.styles["Normal"]
    normal.font.name = "Manrope"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Manrope")
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = docx_rgb(INK_HEX)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.25

    for style_name, size, color in (
        ("Heading 1", 18, BRAND_STRONG_HEX),
        ("Heading 2", 14, BRAND_HEX),
        ("Heading 3", 11.5, BRAND_MID_HEX),
    ):
        style = doc.styles[style_name]
        style.font.name = "Sora"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Sora")
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = docx_rgb(color)
        style.paragraph_format.space_before = Pt(10)
        style.paragraph_format.space_after = Pt(4)


def add_docx_cover(doc: Document, title: str) -> None:
    if LOGO.exists():
        logo_p = doc.add_paragraph()
        logo_p.alignment = WD_ALIGN_PARAGRAPH.LEFT
        logo_run = logo_p.add_run()
        logo_run.add_picture(str(LOGO), width=Inches(2.2))

    eyebrow = doc.add_paragraph()
    eyebrow_run = eyebrow.add_run("RITO Sistemas | Rotinas Inteligentes de Tecnologia e Operação")
    style_docx_run(eyebrow_run, font="IBM Plex Mono", size=8.5, color=ACCENT_HEX)

    title_paragraph = doc.add_paragraph()
    title_run = title_paragraph.add_run(title)
    style_docx_run(title_run, font="Sora", size=21, bold=True, color=BRAND_STRONG_HEX)

    subtitle = doc.add_paragraph()
    subtitle_run = subtitle.add_run("Software sob medida para a rotina da sua empresa")
    style_docx_run(subtitle_run, font="Manrope", size=11, color=INK_SOFT_HEX)

    divider = doc.add_paragraph()
    divider_run = divider.add_run(" ")
    divider_run.font.highlight_color = None


def write_docx_from_markdown(source: Path, output: Path) -> None:
    blocks = parse_markdown(source.read_text(encoding="utf-8"))
    doc = Document()
    setup_docx_document(doc)

    title = document_title(blocks, source.stem)
    add_docx_cover(doc, title)

    first_heading_consumed = False
    for block in blocks:
        kind = block[0]
        if kind == "heading":
            level = block[1]
            text = block[2]
            if not first_heading_consumed:
                first_heading_consumed = True
                continue
            doc.add_heading(text, level=min(level, 3))
        elif kind == "paragraph":
            paragraph = doc.add_paragraph(block[1])
            for run in paragraph.runs:
                style_docx_run(run, font="Manrope", size=10.5, color=INK_HEX)
        elif kind == "bullet":
            paragraph = doc.add_paragraph(block[1], style="List Bullet")
            for run in paragraph.runs:
                style_docx_run(run, font="Manrope", size=10.3, color=INK_HEX)
        elif kind == "number":
            paragraph = doc.add_paragraph(block[1], style="List Number")
            for run in paragraph.runs:
                style_docx_run(run, font="Manrope", size=10.3, color=INK_HEX)
        elif kind == "quote":
            paragraph = doc.add_paragraph()
            run = paragraph.add_run(block[1])
            style_docx_run(run, font="Manrope", size=10.2, italic=True, color=INK_SOFT_HEX)
        elif kind == "code":
            for code_line in block[1].splitlines():
                paragraph = doc.add_paragraph()
                run = paragraph.add_run(code_line)
                style_docx_run(run, font="IBM Plex Mono", size=9, color=BRAND_STRONG_HEX)

    doc.save(output)


def build_pdf_styles():
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="RitoEyebrow",
            parent=styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=8.5,
            leading=11,
            textColor=colors.HexColor(ACCENT_HEX),
            spaceAfter=5,
        )
    )
    styles.add(
        ParagraphStyle(
            name="RitoTitle",
            parent=styles["Title"],
            fontName="Helvetica-Bold",
            fontSize=21,
            leading=24,
            textColor=colors.HexColor(BRAND_STRONG_HEX),
            spaceAfter=8,
        )
    )
    styles.add(
        ParagraphStyle(
            name="RitoHeading2",
            parent=styles["Heading2"],
            fontName="Helvetica-Bold",
            fontSize=13,
            leading=17,
            textColor=colors.HexColor(BRAND_HEX),
            spaceBefore=8,
            spaceAfter=4,
        )
    )
    styles.add(
        ParagraphStyle(
            name="RitoHeading3",
            parent=styles["Heading3"],
            fontName="Helvetica-Bold",
            fontSize=11.2,
            leading=14,
            textColor=colors.HexColor(BRAND_MID_HEX),
            spaceBefore=5,
            spaceAfter=3,
        )
    )
    styles.add(
        ParagraphStyle(
            name="RitoBody",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10.2,
            leading=14.4,
            textColor=colors.HexColor(INK_HEX),
            spaceAfter=5,
        )
    )
    styles.add(
        ParagraphStyle(
            name="RitoBullet",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10.2,
            leading=14.2,
            leftIndent=10,
            textColor=colors.HexColor(INK_HEX),
            spaceAfter=2,
        )
    )
    styles.add(
        ParagraphStyle(
            name="RitoQuote",
            parent=styles["BodyText"],
            fontName="Helvetica-Oblique",
            fontSize=10.2,
            leading=14.2,
            leftIndent=12,
            textColor=colors.HexColor(INK_SOFT_HEX),
            spaceAfter=4,
        )
    )
    return styles


def draw_pdf_frame(pdf_canvas: canvas.Canvas, document_title_text: str) -> None:
    width, height = A4
    pdf_canvas.saveState()
    pdf_canvas.setFillColor(colors.HexColor(BG_LIGHT_HEX))
    pdf_canvas.rect(0, 0, width, height, fill=1, stroke=0)

    pdf_canvas.setFillColor(colors.HexColor(BRAND_STRONG_HEX))
    pdf_canvas.rect(0, height - 16 * mm, width, 16 * mm, fill=1, stroke=0)

    pdf_canvas.setFillColor(colors.HexColor(ACCENT_HEX))
    pdf_canvas.rect(18 * mm, height - 18 * mm, 36 * mm, 1.6 * mm, fill=1, stroke=0)

    if LOGO.exists():
        pdf_canvas.drawImage(
            str(LOGO),
            width - 55 * mm,
            height - 12.5 * mm,
            width=40 * mm,
            preserveAspectRatio=True,
            mask="auto",
        )

    pdf_canvas.setFillColor(colors.HexColor(INK_SOFT_HEX))
    pdf_canvas.setFont("Helvetica", 8.5)
    pdf_canvas.drawString(18 * mm, 10 * mm, "RITO Sistemas")
    pdf_canvas.drawRightString(width - 18 * mm, 10 * mm, f"Página {pdf_canvas.getPageNumber()}")
    pdf_canvas.restoreState()


def write_pdf_from_markdown(source: Path, output: Path) -> None:
    blocks = parse_markdown(source.read_text(encoding="utf-8"))
    title = document_title(blocks, source.stem)
    styles = build_pdf_styles()

    doc = SimpleDocTemplate(
        str(output),
        pagesize=A4,
        rightMargin=18 * mm,
        leftMargin=18 * mm,
        topMargin=28 * mm,
        bottomMargin=16 * mm,
    )

    story = []
    if LOGO.exists():
        story.append(Image(str(LOGO), width=50 * mm, height=18.8 * mm))
        story.append(Spacer(1, 4))

    story.append(Paragraph("RITO Sistemas | Rotinas Inteligentes de Tecnologia e Operação", styles["RitoEyebrow"]))
    story.append(Paragraph(title, styles["RitoTitle"]))
    story.append(Paragraph("Software sob medida para a rotina da sua empresa", styles["RitoBody"]))
    story.append(Spacer(1, 4))

    first_heading_consumed = False
    for block in blocks:
        kind = block[0]
        if kind == "heading":
            level = block[1]
            text = block[2]
            if not first_heading_consumed:
                first_heading_consumed = True
                continue
            story.append(Paragraph(text, styles["RitoHeading2" if level <= 2 else "RitoHeading3"]))
        elif kind == "paragraph":
            story.append(Paragraph(block[1], styles["RitoBody"]))
        elif kind == "bullet":
            story.append(Paragraph(f"• {block[1]}", styles["RitoBullet"]))
        elif kind == "number":
            story.append(Paragraph(block[1], styles["RitoBullet"]))
        elif kind == "quote":
            story.append(Paragraph(block[1], styles["RitoQuote"]))
        elif kind == "code":
            story.append(
                Preformatted(
                    block[1],
                    ParagraphStyle(
                        "CodeBlock",
                        fontName="Courier",
                        fontSize=8.8,
                        leading=11.2,
                        textColor=colors.HexColor(BRAND_STRONG_HEX),
                        backColor=colors.HexColor(BG_HEX),
                        leftIndent=8,
                        rightIndent=8,
                    ),
                )
            )
        story.append(Spacer(1, 1))

    doc.build(
        story,
        onFirstPage=lambda canv, _: draw_pdf_frame(canv, title),
        onLaterPages=lambda canv, _: draw_pdf_frame(canv, title),
    )


def parse_presentation_slides(source: Path):
    slides = []
    current = None
    in_slide_area = False
    skip_objective = False

    for raw_line in source.read_text(encoding="utf-8").splitlines():
        stripped = raw_line.strip()
        match = re.match(r"^###\s+Slide\s+\d+:\s*(.+)$", stripped)
        if match:
            in_slide_area = True
            skip_objective = False
            if current:
                slides.append(current)
            current = {"title": clean_inline(match.group(1)), "lead": "", "details": [], "items": []}
            continue

        if in_slide_area and stripped.startswith("## "):
            if current:
                slides.append(current)
                current = None
            break

        if not current:
            continue

        if not stripped:
            skip_objective = False
            continue

        lower = stripped.lower()
        if lower.startswith("objetivo do slide"):
            skip_objective = True
            continue
        if skip_objective:
            continue
        if lower in LABEL_LINES:
            continue

        bullet_match = re.match(r"^[-*]\s+(.*)$", stripped)
        number_match = re.match(r"^\d+\.\s+(.*)$", stripped)
        if bullet_match:
            current["items"].append(clean_inline(bullet_match.group(1)))
            continue
        if number_match:
            current["items"].append(clean_inline(number_match.group(1)))
            continue

        content = clean_inline(stripped)
        if not current["lead"]:
            current["lead"] = content
        else:
            current["details"].append(content)

    if current:
        slides.append(current)

    return slides


def add_textbox(slide, left: float, top: float, width: float, height: float):
    return slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))


def style_ppt_run(run, *, font: str, size: int, color: str, bold: bool = False) -> None:
    run.font.name = font
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_ppt_textbox(slide, left: float, top: float, width: float, height: float, text: str, *, font: str, size: int, color: str, bold: bool = False, align=PP_ALIGN.LEFT):
    box = add_textbox(slide, left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    style_ppt_run(run, font=font, size=size, color=color, bold=bold)
    return box


def draw_cover_slide(prs: Presentation, slide_data) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(BG_LIGHT_HEX)
    background.line.fill.background()

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(8.8),
        0,
        PptInches(4.533),
        prs.slide_height,
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(BRAND_STRONG_HEX)
    panel.line.fill.background()

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(0.95),
        PptInches(1.08),
        PptInches(1.25),
        PptInches(0.06),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = rgb(ACCENT_HEX)
    accent_bar.line.fill.background()

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), PptInches(0.95), PptInches(0.58), width=PptInches(2.3))
    if MONOGRAM.exists():
        slide.shapes.add_picture(str(MONOGRAM), PptInches(10.25), PptInches(0.72), width=PptInches(1.6))

    add_ppt_textbox(
        slide,
        0.95,
        1.28,
        6.6,
        0.45,
        "RITO Sistemas | Rotinas Inteligentes de Tecnologia e Operação",
        font="IBM Plex Mono",
        size=9,
        color=ACCENT_HEX,
    )

    title = slide_data["lead"] or "RITO Sistemas"
    details = slide_data["details"]
    subtitle = details[0] if details else "Software sob medida para a rotina da sua empresa"
    tertiary = details[1] if len(details) > 1 else "Soluções práticas para micro e pequenas empresas"

    add_ppt_textbox(
        slide,
        0.95,
        2.0,
        6.6,
        0.8,
        title,
        font="Sora",
        size=28,
        color=BRAND_STRONG_HEX,
        bold=True,
    )
    add_ppt_textbox(
        slide,
        0.95,
        2.92,
        6.8,
        1.0,
        subtitle,
        font="Sora",
        size=21,
        color=INK_HEX,
        bold=True,
    )
    add_ppt_textbox(
        slide,
        0.95,
        4.18,
        5.8,
        0.8,
        tertiary,
        font="Manrope",
        size=13,
        color=INK_SOFT_HEX,
    )

    quote_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptInches(9.25),
        PptInches(2.12),
        PptInches(3.15),
        PptInches(2.35),
    )
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = rgb(BRAND_HEX)
    quote_box.line.fill.background()

    add_ppt_textbox(
        slide,
        9.55,
        2.48,
        2.55,
        1.2,
        "Tecnologia prática para organizar, automatizar e dar clareza à operação.",
        font="Manrope",
        size=16,
        color=BG_LIGHT_HEX,
        bold=True,
    )

    add_ppt_textbox(
        slide,
        0.95,
        6.55,
        4.2,
        0.35,
        "Apresentação institucional e comercial",
        font="IBM Plex Mono",
        size=8,
        color=BRAND_MID_HEX,
    )


def split_items(items: Sequence[str]) -> Tuple[List[str], List[str]]:
    if len(items) <= 4:
        return list(items), []
    midpoint = (len(items) + 1) // 2
    return list(items[:midpoint]), list(items[midpoint:])


def draw_regular_slide(prs: Presentation, slide_data, index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(BG_HEX)
    background.line.fill.background()

    rail = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        PptInches(0.58),
        prs.slide_height,
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = rgb(BRAND_STRONG_HEX)
    rail.line.fill.background()

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptInches(4.1),
        PptInches(0.72),
        PptInches(8.45),
        PptInches(5.98),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(SURFACE_HEX)
    card.line.color.rgb = rgb(LINE_HEX)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), PptInches(10.42), PptInches(0.42), width=PptInches(1.95))
    if MONOGRAM.exists():
        slide.shapes.add_picture(str(MONOGRAM), PptInches(0.13), PptInches(6.27), width=PptInches(0.32))

    add_ppt_textbox(
        slide,
        0.95,
        0.58,
        2.6,
        0.28,
        "RITO Sistemas",
        font="IBM Plex Mono",
        size=8,
        color=ACCENT_HEX,
    )
    add_ppt_textbox(
        slide,
        0.95,
        1.0,
        2.65,
        1.55,
        slide_data["title"],
        font="Sora",
        size=24,
        color=BRAND_STRONG_HEX,
        bold=True,
    )

    accent_line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(0.95),
        PptInches(2.55),
        PptInches(0.95),
        PptInches(0.05),
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = rgb(ACCENT_HEX)
    accent_line.line.fill.background()

    lead = slide_data["lead"] or " "
    details = slide_data["details"]
    items = slide_data["items"]

    add_ppt_textbox(
        slide,
        4.48,
        1.1,
        7.55,
        1.15,
        lead,
        font="Sora",
        size=20 if len(lead) > 90 else 22,
        color=BRAND_HEX,
        bold=True,
    )

    top_cursor = 2.2
    for detail in details[:2]:
        add_ppt_textbox(
            slide,
            4.5,
            top_cursor,
            7.3,
            0.62,
            detail,
            font="Manrope",
            size=12,
            color=INK_SOFT_HEX,
        )
        top_cursor += 0.58

    left_items, right_items = split_items(items)
    columns = [left_items] + ([right_items] if right_items else [])
    col_lefts = [4.5, 8.45]

    for column_index, column_items in enumerate(columns):
        y = max(top_cursor + 0.15, 3.15)
        for item in column_items:
            bullet = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                PptInches(col_lefts[column_index]),
                PptInches(y + 0.08),
                PptInches(0.14),
                PptInches(0.14),
            )
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = rgb(ACCENT_HEX)
            bullet.line.fill.background()
            add_ppt_textbox(
                slide,
                col_lefts[column_index] + 0.24,
                y,
                3.35,
                0.55,
                item,
                font="Manrope",
                size=12,
                color=INK_HEX,
            )
            y += 0.55

    add_ppt_textbox(
        slide,
        11.68,
        6.74,
        0.55,
        0.22,
        f"{index:02d}",
        font="IBM Plex Mono",
        size=8,
        color=BRAND_MID_HEX,
        align=PP_ALIGN.RIGHT,
    )


def write_presentation(slides_data, output: Path) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    if slides_data:
        draw_cover_slide(prs, slides_data[0])
    for index, slide_data in enumerate(slides_data[1:], start=2):
        draw_regular_slide(prs, slide_data, index)

    prs.save(output)


def wrap_text(text: str, width: int) -> List[str]:
    words = text.split()
    if not words:
        return [""]

    lines = []
    current = words[0]
    for word in words[1:]:
        if len(current) + 1 + len(word) <= width:
            current += " " + word
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def draw_slide_pdf(source: Path, output: Path) -> None:
    slides = parse_presentation_slides(source)
    page_size = (13.333 * inch, 7.5 * inch)
    pdf_canvas = canvas.Canvas(str(output), pagesize=page_size)
    width, height = page_size

    for index, slide in enumerate(slides, start=1):
        pdf_canvas.setFillColor(colors.HexColor(BG_HEX if index > 1 else BG_LIGHT_HEX))
        pdf_canvas.rect(0, 0, width, height, fill=1, stroke=0)

        if index == 1:
            pdf_canvas.setFillColor(colors.HexColor(BRAND_STRONG_HEX))
            pdf_canvas.rect(width - 4.5 * inch, 0, 4.5 * inch, height, fill=1, stroke=0)
            pdf_canvas.setFillColor(colors.HexColor(ACCENT_HEX))
            pdf_canvas.rect(0.95 * inch, height - 1.18 * inch, 1.25 * inch, 0.05 * inch, fill=1, stroke=0)

            if LOGO.exists():
                pdf_canvas.drawImage(str(LOGO), 0.95 * inch, height - 0.98 * inch, width=2.2 * inch, preserveAspectRatio=True, mask="auto")
            if MONOGRAM.exists():
                pdf_canvas.drawImage(str(MONOGRAM), width - 1.72 * inch, height - 1.58 * inch, width=1.45 * inch, preserveAspectRatio=True, mask="auto")

            details = slide["details"]
            title = slide["lead"] or "RITO Sistemas"
            subtitle = details[0] if details else "Software sob medida para a rotina da sua empresa"
            tertiary = details[1] if len(details) > 1 else "Soluções práticas para micro e pequenas empresas"

            pdf_canvas.setFillColor(colors.HexColor(ACCENT_HEX))
            pdf_canvas.setFont("Helvetica-Bold", 9)
            pdf_canvas.drawString(0.95 * inch, height - 1.38 * inch, "RITO Sistemas | Rotinas Inteligentes de Tecnologia e Operação")

            pdf_canvas.setFillColor(colors.HexColor(BRAND_STRONG_HEX))
            pdf_canvas.setFont("Helvetica-Bold", 28)
            pdf_canvas.drawString(0.95 * inch, height - 2.28 * inch, title)

            pdf_canvas.setFillColor(colors.HexColor(INK_HEX))
            pdf_canvas.setFont("Helvetica-Bold", 20)
            for offset, line in enumerate(wrap_text(subtitle, 38)):
                pdf_canvas.drawString(0.95 * inch, height - (3.0 + offset * 0.34) * inch, line)

            pdf_canvas.setFillColor(colors.HexColor(INK_SOFT_HEX))
            pdf_canvas.setFont("Helvetica", 12)
            for offset, line in enumerate(wrap_text(tertiary, 55)):
                pdf_canvas.drawString(0.95 * inch, height - (4.28 + offset * 0.24) * inch, line)

            pdf_canvas.setFillColor(colors.HexColor(BRAND_HEX))
            pdf_canvas.roundRect(width - 3.78 * inch, height - 4.5 * inch, 3.0 * inch, 1.95 * inch, 10, fill=1, stroke=0)
            pdf_canvas.setFillColor(colors.white)
            pdf_canvas.setFont("Helvetica-Bold", 15)
            quote = "Tecnologia prática para organizar, automatizar e dar clareza à operação."
            y = height - 2.7 * inch
            for line in wrap_text(quote, 27):
                pdf_canvas.drawString(width - 3.48 * inch, y, line)
                y -= 0.28 * inch
        else:
            pdf_canvas.setFillColor(colors.HexColor(BRAND_STRONG_HEX))
            pdf_canvas.rect(0, 0, 0.58 * inch, height, fill=1, stroke=0)
            pdf_canvas.setFillColor(colors.HexColor(SURFACE_HEX))
            pdf_canvas.roundRect(4.1 * inch, 0.72 * inch, 8.45 * inch, 5.98 * inch, 12, fill=1, stroke=0)

            if LOGO.exists():
                pdf_canvas.drawImage(str(LOGO), width - 2.12 * inch, height - 0.77 * inch, width=1.82 * inch, preserveAspectRatio=True, mask="auto")
            if MONOGRAM.exists():
                pdf_canvas.drawImage(str(MONOGRAM), 0.12 * inch, 0.26 * inch, width=0.30 * inch, preserveAspectRatio=True, mask="auto")

            pdf_canvas.setFillColor(colors.HexColor(ACCENT_HEX))
            pdf_canvas.setFont("Helvetica-Bold", 8)
            pdf_canvas.drawString(0.95 * inch, height - 0.72 * inch, "RITO Sistemas")

            pdf_canvas.setFillColor(colors.HexColor(BRAND_STRONG_HEX))
            pdf_canvas.setFont("Helvetica-Bold", 23)
            title_y = height - 1.25 * inch
            for line in wrap_text(slide["title"], 18):
                pdf_canvas.drawString(0.95 * inch, title_y, line)
                title_y -= 0.34 * inch

            pdf_canvas.setFillColor(colors.HexColor(ACCENT_HEX))
            pdf_canvas.rect(0.95 * inch, height - 2.62 * inch, 0.95 * inch, 0.05 * inch, fill=1, stroke=0)

            pdf_canvas.setFillColor(colors.HexColor(BRAND_HEX))
            pdf_canvas.setFont("Helvetica-Bold", 20)
            lead_y = height - 1.5 * inch
            for line in wrap_text(slide["lead"], 46):
                pdf_canvas.drawString(4.48 * inch, lead_y, line)
                lead_y -= 0.3 * inch

            details_y = lead_y - 0.12 * inch
            pdf_canvas.setFillColor(colors.HexColor(INK_SOFT_HEX))
            pdf_canvas.setFont("Helvetica", 11)
            for detail in slide["details"][:2]:
                for line in wrap_text(detail, 65):
                    pdf_canvas.drawString(4.5 * inch, details_y, line)
                    details_y -= 0.21 * inch
                details_y -= 0.06 * inch

            left_items, right_items = split_items(slide["items"])
            for column_index, column_items in enumerate((left_items, right_items)):
                if not column_items:
                    continue
                x = 4.5 * inch if column_index == 0 else 8.45 * inch
                y = min(details_y - 0.15 * inch, height - 3.12 * inch)
                for item in column_items:
                    pdf_canvas.setFillColor(colors.HexColor(ACCENT_HEX))
                    pdf_canvas.circle(x + 0.05 * inch, y - 0.03 * inch, 0.05 * inch, fill=1, stroke=0)
                    pdf_canvas.setFillColor(colors.HexColor(INK_HEX))
                    pdf_canvas.setFont("Helvetica", 11.5)
                    first_line = True
                    for line in wrap_text(item, 28):
                        prefix = "" if not first_line else ""
                        pdf_canvas.drawString(x + 0.18 * inch, y, f"{prefix}{line}")
                        y -= 0.18 * inch
                        first_line = False
                    y -= 0.14 * inch

            pdf_canvas.setFillColor(colors.HexColor(BRAND_MID_HEX))
            pdf_canvas.setFont("Helvetica-Bold", 8)
            pdf_canvas.drawRightString(width - 0.48 * inch, 0.3 * inch, f"{index:02d}")

        pdf_canvas.showPage()

    pdf_canvas.save()


def create_deliverables_readme() -> None:
    readme = OUT_ROOT / "README.md"
    readme.write_text(
        "\n".join(
            [
                "# RITO Business Kit",
                "",
                "Esta pasta reúne os entregáveis finais da RITO Sistemas em formatos prontos para uso comercial e institucional.",
                "",
                "## Estrutura",
                "",
                "- `editable-docx/`: arquivos editáveis em Word para briefing, proposta, orçamento, contrato, follow-up e materiais operacionais.",
                "- `presentation/`: apresentação institucional em PowerPoint.",
                "- `pdf/`: materiais prontos para distribuição em PDF.",
                "",
                "## Arquivos principais",
                "",
                "- `editable-docx/rito-contract-template.docx`",
                "- `editable-docx/rito-commercial-proposal-template.docx`",
                "- `editable-docx/rito-estimate-template.docx`",
                "- `editable-docx/rito-client-briefing-form.docx`",
                "- `editable-docx/rito-pricing-guide.docx`",
                "- `pdf/rito-branding-guide.pdf`",
                "- `pdf/rito-software-brand-system.pdf`",
                "- `presentation/rito-company-presentation.pptx`",
                "- `pdf/rito-company-presentation.pdf`",
                "- `pdf/rito-company-onepager.pdf`",
                "",
                "## Como regenerar",
                "",
                "```bash",
                "python3 scripts/generate_final_business_assets.py",
                "```",
                "",
                "## Observação",
                "",
                "Os arquivos desta pasta são gerados a partir da base em `docs/sales/`, `docs/legal/` e `docs/brand/`. Quando os documentos-fonte forem refinados, execute o script novamente para atualizar o kit final.",
                "",
            ]
        ),
        encoding="utf-8",
    )


def main() -> None:
    ensure_dirs()
    create_deliverables_readme()

    created = []

    for filename, source in existing_sources(DOCX_SOURCES):
        output = DOCX_DIR / filename
        write_docx_from_markdown(source, output)
        created.append(output)

    for filename, source in existing_sources(PDF_SOURCES):
        output = PDF_DIR / filename
        write_pdf_from_markdown(source, output)
        created.append(output)

    if PRESENTATION_MD.exists():
        slides = parse_presentation_slides(PRESENTATION_MD)
        if slides:
            write_presentation(slides, PRESENTATION_PPTX)
            draw_slide_pdf(PRESENTATION_MD, PRESENTATION_PDF)
            created.extend([PRESENTATION_PPTX, PRESENTATION_PDF])

    for path in created:
        print(path.relative_to(ROOT))


if __name__ == "__main__":
    main()
